
import os
import requests
import io 
import json
import re
import base64
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

try:
    from groq import Groq
except Exception:
    Groq = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import docx2txt
except Exception:
    docx2txt = None

try:
    from google import genai
    from google.genai import types
except Exception:
    genai = None
    types = None

from database import get_db_connection

app = FastAPI()


# 👇 Yeh CORS Middleware setup add karein
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Har frontend (Live Server + Vercel) ko allow karega
    allow_credentials=True,
    allow_methods=["*"],  # All methods (GET, POST, etc.)
    allow_headers=["*"],
)

@app.get("/health")
@app.get("/api/health")
async def health_check():
    return {"status": "online", "message": "Backend server is live and working!"}


# ====================================================
# CRITICAL: API KEYS CONFIGURATION (LOAD FROM ENV)
# ====================================================
GROQ_KEY = os.getenv("GROQ_KEY", "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")      
OPENROUTER_KEY = os.getenv("OPENROUTER_KEY", "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")  
GEMINI_KEY = os.getenv("GEMINI_KEY", "AQ.Ab8RN6K5igFNFB0ayrDT3fELaPbHUh0eOeZI75jx1-CG2f5AvA")

class ChatRequest(BaseModel):
    message: str
    email: str = "guest@gmail.com"


# =====================================================================
# ENGINE 1: DIRECT GOOGLE GEMINI (PRIMARY ENGINE FOR BOTH TEXT & VISION)
# =====================================================================
def call_direct_gemini_api(prompt_text: str, data_url: str = "") -> str:
    try:
        if not GEMINI_KEY:
            return "ERROR"
            
        client = genai.Client(api_key=GEMINI_KEY)
        
        # Agar image present hai
        if data_url.strip():
            if "," in data_url:
                raw_b64 = data_url.split(",")[1]
            else:
                raw_b64 = data_url

            image_bytes = base64.b64decode(raw_b64)
            image_part = types.Part.from_bytes(
                data=image_bytes,
                mime_type="image/jpeg"
            )
            contents_payload = [prompt_text, image_part]
        else:
            contents_payload = prompt_text

        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=contents_payload
        )
        if response and response.text:
            return response.text
    except Exception as e:
        print(f"Direct Gemini Engine Connection Log: {e}")
    return "ERROR"


# =====================================================================
# ENGINE 2: OPENROUTER (SECONDARY FALLBACK ENGINE)
# =====================================================================
def call_openrouter_api(messages: list) -> str:
    try:
        if not OPENROUTER_KEY:
            return "ERROR"

        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "google/gemini-2.5-flash",
                "messages": messages
            },
            timeout=10
        )
        if response.status_code == 200:
            return response.json()['choices'][0]['message']['content']
    except Exception as e:
        print(f"OpenRouter Connection Log: {e}")
    return "ERROR"


# =====================================================================
# ENGINE 3: GROQ (TERTIARY ENGINE - ACTIVE MODEL)
# =====================================================================
def call_groq_api(prompt_text: str) -> str:
    try:
        if not GROQ_KEY:
            return "ERROR"

        client = Groq(api_key=GROQ_KEY)
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt_text}],
            temperature=0.5
        )
        return completion.choices[0].message.content
    except Exception as e:
        print(f"Groq Connection Log: {e}")
    return "ERROR"


# =====================================================================
# ENGINE 4: FREE FALLBACK VISION ENGINE (QUATERNARY ENGINE)
# =====================================================================
def call_free_vision_fallback(data_url: str, prompt_text: str) -> str:
    try:
        payload = {
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt_text},
                        {"type": "image_url", "image_url": {"url": data_url}}
                    ]
                }
            ],
            "model": "openai"
        }
        res = requests.post("https://text.pollinations.ai/", json=payload, timeout=12)
        if res.status_code == 200 and res.text.strip():
            return res.text
    except Exception as e:
        print(f"Free Fallback Engine Log: {e}")
    return "ERROR"


# =====================================================================
# CORE INTEGRATED MULTI-FALLBACK ENGINE
# =====================================================================
def get_fallback_ai_response(messages: list, raw_b64_image: str = "", prompt_text: str = "") -> str:
    if not prompt_text:
        prompt_text = "\n\n".join([f"{m.get('role', 'user')}: {m.get('content', '')}" for m in messages])

    # 1. Primary Engine: Direct Google Gemini API
    print("Trying Engine 1: Direct Google Gemini...")
    res = call_direct_gemini_api(prompt_text=prompt_text, data_url=raw_b64_image)
    if res != "ERROR":
        return res

    # 2. Secondary Engine: OpenRouter
    print("Shifting to Engine 2: OpenRouter...")
    res = call_openrouter_api(messages)
    if res != "ERROR":
        return res

    # 3. Tertiary Engine: Groq
    print("Shifting to Engine 3: Groq...")
    res = call_groq_api(prompt_text)
    if res != "ERROR":
        return res

    # 4. Quaternary Engine: Open Fallback Vision (If Image present)
    if raw_b64_image:
        print("Shifting to Engine 4: Open Fallback Vision Engine...")
        res = call_free_vision_fallback(raw_b64_image, prompt_text)
        if res != "ERROR":
            return res

    return "Tamam AI Engines respond nahi kar rahe. Meharbani karke backend terminal logs aur API Keys check karein."


# AI STUDY CHAT
@app.post("/ai/study-chat")
async def study_chat(payload: ChatRequest):
    system_prompt = (
        "You are an expert, friendly AI Study Assistant. Provide highly structured, "
        "beautiful, and easy-to-read responses using clear formatting rules.\n"
        "Rules:\n"
        "1. Use '### Heading Name' for major topics or sub-sections.\n"
        "2. Use single asterisk bullets '* Keypoint: details' for bullet items.\n"
        "3. Use '**text**' to bold critical key terms.\n"
        "4. Always add a newline character between paragraphs and headings to avoid text crowding.\n"
        "Prioritize scannability and structural hierarchy."
    )
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": payload.message}
    ]
    
    ai_res = get_fallback_ai_response(messages, prompt_text=f"{system_prompt}\n\nUser Question: {payload.message}")
    
    db = get_db_connection()
    if db:
        try:
            cursor = db.cursor(dictionary=True)
            cursor.execute("SELECT user_id FROM users WHERE email = %s", (payload.email,))
            user_data = cursor.fetchone()
            
            user_id = 4 
            if user_data:
                user_id = user_data['user_id']
                
            query = """
                INSERT INTO chat_history (user_id, prompt, response)
                VALUES (%s, %s, %s)
            """
            cursor.execute(query, (user_id, payload.message, ai_res))
            db.commit() 
            cursor.close()
            print("Chat history successfully saved to database!")
        except Exception as db_err:
            if db:
                db.rollback()
            print(f"Database history log failed: {db_err}")
        finally:
            db.close()

    return {"response": ai_res}


# AI SUMMARIZE
@app.post("/ai/summarize")
async def summarize(file: UploadFile = File(...), email: str = Form("guest@gmail.com")):
    extracted_text = ""
    filename = file.filename.lower()
    
    try:
        file_bytes = await file.read()
        
        if filename.endswith('.txt'):
            extracted_text = file_bytes.decode("utf-8")
        elif filename.endswith('.pdf'):
            pdf_file = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_file)
            for page in reader.pages:
                text = page.extract_text()
                if text: extracted_text += text + "\n"
        elif filename.endswith('.ppt') or filename.endswith('.pptx'):
            ppt_file = io.BytesIO(file_bytes)
            prs = Presentation(ppt_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        extracted_text += shape.text + "\n"
        else:
            raise HTTPException(status_code=400, detail="Sirf .txt, .pdf, aur .ppt files allowed hain.")
            
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="File khali hai ya is me se text read nahi ho saka.")
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"File reading error: {str(e)}")

    system_prompt = (
        "You are an expert academic summarizer. Summarize the provided text in very "
        "simple, easy-to-understand words. Use clear headings or bullet points if necessary."
    )
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"Please summarize this text: {extracted_text}"}
    ]
    
    ai_res = get_fallback_ai_response(messages, prompt_text=f"{system_prompt}\n\nPlease summarize this text:\n{extracted_text}")
    
    db = get_db_connection()
    if db:
        try:
            cursor = db.cursor(dictionary=True)
            cursor.execute("SELECT user_id FROM users WHERE email = %s", (email,))
            user_data = cursor.fetchone()
            user_id = 4 if not user_data else user_data['user_id']
            
            query = """
                INSERT INTO documents (user_id, file_name, summary_text)
                VALUES (%s, %s, %s)
            """
            cursor.execute(query, (user_id, file.filename, ai_res))
            db.commit()
            cursor.close()
        except Exception as db_err:
            print(f"Database document log failed: {db_err}")
        finally:
            db.close()

    return {"response": ai_res}


# AI QUIZ GENERATOR 
class QuizRequest(BaseModel):
    topic: str
    email: str = "guest@gmail.com"

@app.post("/ai/generate-quiz")
async def generate_quiz(payload: QuizRequest):
    system_prompt = (
        "You are an expert quiz master. Generate a quiz with exactly 5 multiple choice questions "
        "about the requested topic. Your entire response must be a single valid JSON list, "
        "with absolutely no markdown formatting, no code blocks (like ```json), and no extra text. "
        "Each object in the list must match this exact format: "
        '{"question": "Question text here", "a": "Option A", "b": "Option B", "c": "Option C", "d": "Option D", "answer": "a"}'
    )
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"Generate a quiz about: {payload.topic}"}
    ]
    
    full_prompt = f"{system_prompt}\n\nGenerate a quiz about: {payload.topic}"
    ai_raw_res = get_fallback_ai_response(messages, prompt_text=full_prompt)
    
    clean_json_str = re.sub(r"```json|```", "", ai_raw_res).strip()
    
    try:
        questions_list = json.loads(clean_json_str)
    except Exception as parse_err:
        print(f"JSON Parsing failed, raw response: {ai_raw_res}")
        raise HTTPException(status_code=500, detail="AI response pattern mismatch. Please try again.")

    db = get_db_connection()
    quiz_id = None
    
    if db:
        try:
            cursor = db.cursor(dictionary=True)
            cursor.execute("SELECT user_id FROM users WHERE email = %s", (payload.email,))
            user_data = cursor.fetchone()
            user_id = 4 if not user_data else user_data['user_id']
            
            quiz_query = """
                INSERT INTO quizzes (user_id, topic, score, total_questions)
                VALUES (%s, %s, %s, %s)
            """
            cursor.execute(quiz_query, (user_id, payload.topic, 0, len(questions_list)))
            db.commit()
            
            quiz_id = cursor.lastrowid
            
            questions_query = """
                INSERT INTO quiz_questions (quiz_id, question_text, option_a, option_b, option_c, option_d, correct_option)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
            """
            
            for q in questions_list:
                cursor.execute(questions_query, (
                    quiz_id,
                    q.get('question'),
                    q.get('a'),
                    q.get('b'),
                    q.get('c'),
                    q.get('d'),
                    q.get('answer', 'a').lower()
                ))
            db.commit()
            cursor.close()
        except Exception as db_err:
            print(f"Database Quiz log failed: {db_err}")
        finally:
            db.close()

    return {"quiz_id": quiz_id, "questions": questions_list}


# MULTI-UPLOAD DOCUMENT EXPLAINER CODE (FIXED DB SAVE)
@app.post("/api/multi-upload-explain")
async def multi_upload_explain(
    file: UploadFile = File(...), 
    email: str = Form("guest@gmail.com")
):
    system_prompt = (
        "You are an advanced AI Academic Assistant specializing in detailed document analysis.\n"
        "Your task is to thoroughly analyze the provided document content and explain it in deep detail using simple, clear, and highly organized English prose.\n\n"
        "Formatting Rules:\n"
        "1. Use '### Heading Name' for major topics or sub-sections.\n"
        "2. Use single asterisk bullets '* Keypoint: details' for bullet items.\n"
        "3. Use '**text**' to bold critical key terms.\n"
        "4. Always add a newline character between paragraphs and headings to avoid text crowding.\n"
        "Prioritize clean structural hierarchy and scannability."
    )
    
    extracted_text = ""
    filename = file.filename.lower()
    
    # 1. File Reading & Processing
    try:
        file_bytes = await file.read()
        
        if filename.endswith('.docx'):
            docx_file = io.BytesIO(file_bytes)
            extracted_text = docx2txt.process(docx_file)
        elif filename.endswith('.pdf'):
            pdf_file = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_file)
            for page in reader.pages:
                text = page.extract_text()
                if text: 
                    extracted_text += text + "\n"
        elif filename.endswith('.txt'):
            extracted_text = file_bytes.decode("utf-8")
        elif filename.endswith('.ppt') or filename.endswith('.pptx'):
            ppt_file = io.BytesIO(file_bytes)
            prs = Presentation(ppt_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        extracted_text += shape.text + "\n"
        else:
            raise HTTPException(
                status_code=400, 
                detail="This file type is not supported. Please upload .docx, .pdf, .txt, or .ppt files."
            )
            
        if not extracted_text.strip():
            raise HTTPException(
                status_code=400, 
                detail="We couldn't extract any readable text from this file."
            )
            
    except HTTPException as http_ex:
        raise http_ex
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"File reading error: {str(e)}")

    # 2. AI Processing
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"Analyze and explain the following extracted text in highly detailed, plain English:\n\n{extracted_text}"}
    ]
    
    try:
        full_p = f"{system_prompt}\n\nAnalyze and explain this text:\n{extracted_text}"
        ai_explanation = get_fallback_ai_response(messages, prompt_text=full_p)
    except Exception as ai_err:
        raise HTTPException(status_code=500, detail=f"AI Generation Failed: {str(ai_err)}")
    
  # 3. FIXED & VERIFIED DATABASE LOGGING SECTION
    db = None
    try:
        db = get_db_connection()
        if not db:
            print("❌ DATABASE ERROR: get_db_connection() returned None! Check DB Credentials.")
        else:
            cursor = db.cursor(dictionary=True) if hasattr(db, 'cursor') else db.cursor()
            
            # A. PRINT ACTUAL DATABASE NAME (Check karein GUI mein yehi open hai ya nahi)
            try:
                cursor.execute("SELECT DATABASE() as db_name")
                current_db = cursor.fetchone()
                db_name_str = current_db.get('db_name') if isinstance(current_db, dict) else current_db[0]
                print(f"📌 ACTIVE DATABASE NAME: [{db_name_str}]")
            except Exception:
                pass

            # B. SAFE USER FETCHING (Avoid Foreign Key Failure)
            cursor.execute("SELECT user_id FROM users WHERE email = %s", (email,))
            user_data = cursor.fetchone()
            
            user_id = None
            if user_data:
                user_id = user_data.get('user_id') if isinstance(user_data, dict) else user_data[0]
            else:
                # Fallback: Pick the very first existing user from DB to prevent foreign key error
                cursor.execute("SELECT user_id FROM users LIMIT 1")
                first_user = cursor.fetchone()
                if first_user:
                    user_id = first_user.get('user_id') if isinstance(first_user, dict) else first_user[0]
                else:
                    user_id = 1 # Absolute fallback

            # C. INSERT INTO TABLE
            insert_query = """
                INSERT INTO multi_uploaded_docs (user_id, file_name, explanation_text)
                VALUES (%s, %s, %s)
            """
            cursor.execute(insert_query, (user_id, file.filename, ai_explanation))
            
            # FORCE COMMIT TO HARD DISK
            db.commit()
            
            # D. LIVE VERIFICATION READ
            cursor.execute("SELECT id, file_name FROM multi_uploaded_docs ORDER BY id DESC LIMIT 1")
            saved_row = cursor.fetchone()
            print(f"✅ SUCCESSFULLY SAVED & VERIFIED IN DB: {saved_row}")

            cursor.close()

    except Exception as db_err:
        print(f"❌ DATABASE INSERT ERROR: {str(db_err)}")
        if db:
            try:
                db.rollback()
            except Exception:
                pass
    finally:
        if db:
            try:
                db.close()
            except Exception:
                pass

    return {"explanation": ai_explanation}


# CODE EXPLAINER ENDPOINT

class CodeExplanationRequest(BaseModel):
    code: str
    language: str = "auto"
    target_language: str = "Roman Urdu/Hindi"
    email: str = "guest@gmail.com"


@app.post("/api/explain-code")
async def explain_code(data: CodeExplanationRequest):
    if not data.code.strip():
        raise HTTPException(status_code=400, detail="Code content cannot be empty.")

    if data.target_language == "Roman Urdu/Hindi":
        language_instruction = (
            "IMPORTANT LANGUAGE RULE: Write the ENTIRE explanation in friendly, natural 'Roman Urdu/Hindi' "
            "(conversational language used in WhatsApp/casual messages, e.g., 'Yeh function array ko loop karta hai...').\n"
            "Keep technical programming keywords (like function, variable, loop, class, array, return, scope) in exact English, "
            "but explain the reasoning and context in Roman Urdu."
        )
    elif data.target_language == "Urdu Script":
        language_instruction = (
            "IMPORTANT LANGUAGE RULE: Write the ENTIRE explanation in standard, fluent Urdu Script (اردو رسم الخط). "
            "Keep core technical terms in brackets where helpful for computer science students."
        )
    else:
        language_instruction = (
            "IMPORTANT LANGUAGE RULE: Write the ENTIRE explanation in professional, clear, and educational English."
        )

    system_prompt = (
        "You are a World-Class Computer Science Professor and Expert Senior Developer.\n"
        "Your primary goal is to explain code in EXTREME DEPTH so that even beginner CS students can grasp every single detail.\n\n"
        f"{language_instruction}\n\n"
        "EXPLANATION STRUCTURE REQUIREMENT (Follow these EXACT sections in order):\n"
        "### 🎯 Executive Summary & Purpose\n"
        "- Explain overall kya kaam ho raha hai aur real software development mein yeh code kyun use hota hai.\n\n"
        "### 🔬 Detailed Line-by-Line / Block Analysis\n"
        "- Code ke HAR ek block, statement, variable, syntax, aur property ko deeply explain karein.\n"
        "- Agar koi special keyword (jaise `margin: 0`, `const`, `map()`, `border-box`) use hua hai to batayein ke woh exactly kya control karta hai aur background mein kaise kaam karta hai.\n\n"
        "### 💡 Real-World Visual Analogy\n"
        "- CS students ke liye ek aasan, everyday life se visual example/analogy dein taaki concept hamesha ke liye clear ho jaye.\n\n"
        "### ⚠️ Common Beginner Mistakes & Best Practices\n"
        "- Batayein ke students yahan kya galti karte hain (e.g., forgetting semi-colons, scope issues, wrong syntax) aur iska best modern alternative kya hai.\n\n"
        "Formatting Rules:\n"
        "1. Use '### Heading Name' for major topics.\n"
        "2. Use single asterisk bullets '* Highlight: explanation' for points.\n"
        "3. Wrap code keywords, variables, or properties in backticks (e.g. `margin`, `padding`, `display`).\n"
        "4. Always add double line breaks between paragraphs and sections to keep text breathable.\n"
    )
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"Programming Language: {data.language}\n\nSource Code to Explain:\n```\n{data.code}\n```"}
    ]
    
    try:
        full_p = f"{system_prompt}\n\nProgramming Language: {data.language}\n\nCode:\n{data.code}"
        ai_explanation = get_fallback_ai_response(messages, prompt_text=full_p)
    except Exception as ai_err:
        raise HTTPException(status_code=500, detail=f"AI Code Analysis Failed: {str(ai_err)}")
    
    db = get_db_connection()
    if db:
        try:
            cursor = db.cursor(dictionary=True)
            cursor.execute("SELECT user_id FROM users WHERE email = %s", (data.email,))
            user_data = cursor.fetchone()
            user_id = 4 if not user_data else user_data['user_id']
            
            query = """
                INSERT INTO code_explanations (user_id, programming_language, raw_code, explanation_text)
                VALUES (%s, %s, %s, %s)
            """
            logged_lang = f"{data.language} ({data.target_language})"
            cursor.execute(query, (user_id, logged_lang, data.code, ai_explanation))
            db.commit()
            cursor.close()
        except Exception as db_err:
            print(f"Database code log failed: {db_err}")
        finally:
            db.close()

    return {"explanation": ai_explanation}



#IMAGE TO TEXT CODE:

class ImageExplanationRequest(BaseModel):
    image_base64: str
    target_language: str
    custom_prompt: str = ""
    email: str = "guest@gmail.com"

# API ENDPOINT: EXPLAIN IMAGE / DIAGRAM
@app.post("/api/explain-image")
async def explain_image(data: ImageExplanationRequest):
    if not data.image_base64.strip():
        raise HTTPException(status_code=400, detail="Image content cannot be empty.")

    # 1. Language Rules
    if data.target_language == "Roman Urdu/Hindi":
        language_instruction = (
            "IMPORTANT LANGUAGE RULE: Write the entire explanation in friendly, natural 'Roman Urdu/Hindi' "
            "(conversational language used in casual chat/WhatsApp, e.g., 'Yeh diagram 3 components ko show kar raha hai...').\n"
            "Keep technical terms (like Class, Database, Gateway, Flowchart, Flow, Step, Variable, Loop) in English text."
        )
    elif data.target_language == "Urdu Script":
        language_instruction = (
            "IMPORTANT LANGUAGE RULE: Write the entire explanation in standard, fluent Urdu Script (اردو رسم الخط). "
            "Keep technical software/diagram terms in readable brackets where needed."
        )
    else:
        language_instruction = (
            "IMPORTANT LANGUAGE RULE: Write the entire explanation in clear, highly professional academic English."
        )

    system_prompt = (
        "You are an Expert AI Computer Science Professor and Multimodal Vision Specialist.\n"
        "Your task is to analyze the provided image (diagram, flowchart, UML, circuit, ERD, architecture, or handwritten notes) in EXTREME DEPTH for students.\n\n"
        f"{language_instruction}\n\n"
        "REQUIRED EXPLANATION STRUCTURE (Strictly follow these section headers):\n"
        "### 🎯 Overview & Context\n"
        "- Identify what this image/diagram represents and its main purpose.\n\n"
        "### 🔬 Detailed Step-by-Step Breakdown\n"
        "- Explain every component, shape, arrow, logic block, or written code line present in the image.\n"
        "- Detail how data/process flows from start to finish.\n\n"
        "### 💡 Real-World Visual Analogy\n"
        "- Give a simple, relatable real-world example/analogy so students can easily visualize the concept.\n\n"
        "### ⚠️ Key Takeaways & Exam Points\n"
        "- Highlight important takeaways or common mistakes students should avoid in exams regarding this diagram/concept.\n\n"
        "Formatting Rules:\n"
        "1. Use '### Heading Name' for major sections.\n"
        "2. Use single asterisk bullets '* Highlight: detail' for points.\n"
        "3. Add double line breaks between paragraphs for clean readability.\n"
    )

    # 2. Base64 cleanup and Mime Type identification
    raw_b64 = data.image_base64
    mime_type = "image/jpeg"
    
    if "," in raw_b64:
        header, raw_b64 = raw_b64.split(",", 1)
        if "data:" in header and ";base64" in header:
            mime_type = header.split(";")[0].replace("data:", "")

    user_text = data.custom_prompt.strip() if data.custom_prompt.strip() else "Please explain this image/diagram in detail."
    full_prompt_text = f"{system_prompt}\n\n{user_text}"

    # 3. Call Gemini Vision Model using your GEMINI_KEY
    try:
        # Client initialize using your GEMINI_KEY variable
        client = genai.Client(api_key=GEMINI_KEY)
        
        image_bytes = base64.b64decode(raw_b64)
        
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                types.Part.from_bytes(
                    data=image_bytes,
                    mime_type=mime_type
                ),
                full_prompt_text
            ]
        )
        ai_explanation = response.text
    except Exception as ai_err:
        print(f"Gemini API Call Failed: {ai_err}")
        raise HTTPException(status_code=500, detail=f"Vision AI Analysis Failed: {str(ai_err)}")

    # 4. Save result to database
    db = get_db_connection()
    if db:
        try:
            cursor = db.cursor(dictionary=True)
            cursor.execute("SELECT user_id FROM users WHERE email = %s", (data.email,))
            user_data = cursor.fetchone()
            user_id = user_data['user_id'] if user_data else 4

            query = """
                INSERT INTO image_explanations (user_id, language, custom_prompt, explanation_text)
                VALUES (%s, %s, %s, %s)
            """
            cursor.execute(query, (user_id, data.target_language, data.custom_prompt, ai_explanation))
            db.commit()
            cursor.close()
        except Exception as db_err:
            print(f"Database image log error: {db_err}")
        finally:
            db.close()

    return {"explanation": ai_explanation}




# SIGNUP ENDPOINT
class SignupRequest(BaseModel):
    username: str
    email: str
    password: str

class LoginRequest(BaseModel):
    email: str
    password: str

@app.post("/auth/signup")
async def signup(payload: SignupRequest):
    db = get_db_connection()
    if not db:
        raise HTTPException(status_code=500, detail="Database connection failure")
    try:
        cursor = db.cursor(dictionary=True)
        cursor.execute("SELECT user_id FROM users WHERE email = %s", (payload.email,))
        if cursor.fetchone():
            raise HTTPException(status_code=400, detail="Email pehle se registered hai!")
        
        query = "INSERT INTO users (username, email, password_hash) VALUES (%s, %s, %s)"
        cursor.execute(query, (payload.username, payload.email, payload.password))
        db.commit()
        cursor.close()
        return {"message": "Account successfully created!"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        db.close()

# LOGIN ENDPOINT
@app.post("/auth/login")
async def login(payload: LoginRequest):
    db = get_db_connection()
    if not db:
        raise HTTPException(status_code=500, detail="Database connection failure")
    try:
        cursor = db.cursor(dictionary=True)
        query = "SELECT user_id, username, email FROM users WHERE email = %s AND password_hash = %s"
        cursor.execute(query, (payload.email, payload.password))
        user = cursor.fetchone()
        cursor.close()
        
        if not user:
            raise HTTPException(status_code=401, detail="Ghalat email ya password!")
            
        return {
            "token": "fake-jwt-token",
            "email": user['email'],
            "username": user['username']
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        db.close()



if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)